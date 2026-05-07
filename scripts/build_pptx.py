#!/usr/bin/env python3
"""Build Sentinel Ops — Clean 9-slide deck (5-second readable per slide)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Dark theme palette ────────────────────────────────────────────────────────
BG       = RGBColor(0x0B, 0x12, 0x20)
SURFACE  = RGBColor(0x14, 0x1E, 0x2E)
TEXT     = RGBColor(0xF8, 0xFA, 0xFC)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
ACCENT   = RGBColor(0x22, 0xC5, 0x5E)   # green
ACCENT2  = RGBColor(0x38, 0xBD, 0xF8)   # blue
TEAL     = RGBColor(0x00, 0xB8, 0x8A)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def bg(slide): fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = BG

def tx(slide, s, l, t, w, h, sz=18, bold=False, col=TEXT, ali=PP_ALIGN.LEFT, italic=False):
    b = slide.shapes.add_textbox(l, t, w, h)
    b.word_wrap = True
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = ali
    r = p.add_run(); r.text = s
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = col

def bul(slide, items, l, t, w, h, sz=15, col=TEXT, sp=8):
    b = slide.shapes.add_textbox(l, t, w, h)
    b.word_wrap = True; tf = b.text_frame; tf.word_wrap = True
    first = True
    for i in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_before = Pt(sp)
        r = p.add_run(); r.text = f"• {i}"
        r.font.size = Pt(sz); r.font.color.rgb = col

def pill(slide, s, l, t, w, h, bg=ACCENT, sz=14):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = bg; sh.line.fill.background()
    tf = sh.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = s; r.font.size = Pt(sz); r.font.bold = True; r.font.color.rgb = WHITE

def divider(slide, l, t, w, col=ACCENT, pt=Pt(3)):
    sh = slide.shapes.add_shape(1, l, t, w, pt)
    sh.fill.solid(); sh.fill.fore_color.rgb = col; sh.line.fill.background()

def slide_title(slide, main, sub=None):
    tx(slide, main, Inches(0.6), Inches(0.28), Inches(12), Inches(0.8),
       sz=32, bold=True, col=TEXT)
    if sub:
        tx(slide, sub, Inches(0.6), Inches(0.95), Inches(12), Inches(0.4),
           sz=14, col=MUTED)
    divider(slide, Inches(0.6), Inches(1.32), Inches(12.1), col=ACCENT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
tx(s, "Sentinel Ops",   Inches(0.7), Inches(1.4),  Inches(12), Inches(1.3), sz=72, bold=True, col=WHITE)
tx(s, "Lean Managed Cybersecurity for Malaysian SMEs",  Inches(0.7), Inches(2.8), Inches(12), Inches(0.6), sz=26, col=ACCENT)
tx(s, "Monitoring  •  Hardening  •  Audit-Ready Reporting", Inches(0.7), Inches(3.45), Inches(12), Inches(0.45), sz=16, col=MUTED)
tx(s, "without building a full SOC.", Inches(0.7), Inches(3.9), Inches(12), Inches(0.45), sz=16, col=MUTED)
divider(s, Inches(0.7), Inches(4.5), Inches(3.5), col=ACCENT2)
tx(s, "[Your Name]",    Inches(0.7), Inches(4.8),  Inches(12), Inches(0.5), sz=20, bold=True, col=WHITE)
tx(s, "Founder / Operator, Sentinel Ops", Inches(0.7), Inches(5.25), Inches(12), Inches(0.4), sz=14, col=MUTED)
tx(s, "[email]   [WhatsApp]   [LinkedIn]", Inches(0.7), Inches(5.65), Inches(12), Inches(0.35), sz=13, col=RGBColor(0x60,0x80,0x90))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM  (5-second: noise, no time, expensive SOC, compliance)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
slide_title(s, "The Problem", "Why growing companies struggle with cybersecurity")

pills_data = [
    ("Noise, not signal",      "Security tools fire alerts.\nNobody triages them.",     ACCENT),
    ("No time",                 "Small IT teams are fighting\nfires, not security.",       ACCENT2),
    ("Inconsistent hygiene",    "Hardening falls off the\npriority list.",               TEAL),
    ("Compliance rising",        "PDPA + enterprise questionnaires\n= real pressure.",     RGBColor(0xA7,0x8A,0xF5)),
]
pw = Inches(2.8); pg = Inches(0.38)
total = 4*pw + 3*pg; sx = (prs.slide_width - total) / 2
for i, (title, desc, col) in enumerate(pills_data):
    x = sx + i*(pw+pg)
    pill(s, title, x, Inches(1.65), pw, Inches(0.45), bg=col, sz=12)
    tx(s, desc, x, Inches(2.2), pw, Inches(1.4), sz=14, col=TEXT)

# Bottom callout — no paragraphs
box = s.shapes.add_shape(1, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.2))
box.fill.solid(); box.fill.fore_color.rgb = SURFACE; box.line.color.rgb = ACCENT; box.line.width = Pt(2)
tx(s, "Most companies need better security operations —", Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.45), sz=16, bold=True, col=WHITE)
tx(s, "but not the cost and complexity of a full SOC.",   Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.45), sz=16, bold=True, col=ACCENT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OUR ANSWER  (5-second: what we do in 5 bullet-pills)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
slide_title(s, "What Sentinel Ops Does", "The lean managed security layer for critical systems")

offers = [
    ("Managed Monitoring",    "24/7 coverage for Linux / cloud. Real alerts, not noise."),
    ("Alert Triage",           "We filter noise and surface what needs attention."),
    ("Baseline Hardening",     "Practical CIS-aligned controls on approved systems."),
    ("Monthly Reporting",      "Management-ready — findings, risk trends, next steps."),
    ("Fast Pilot",             "30-day limited-scope pilot. Prove value, no overcommitment."),
]
for i, (title, desc) in enumerate(offers):
    y = Inches(1.55) + i*Inches(0.98)
    pill(s, title, Inches(0.6), y, Inches(2.4), Inches(0.4), bg=ACCENT, sz=12)
    tx(s, desc, Inches(3.15), y+Inches(0.05), Inches(9.7), Inches(0.4), sz=14, col=TEXT)

box = s.shapes.add_shape(1, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.65))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x0D,0x28,0x1A); box.line.color.rgb = ACCENT; box.line.width = Pt(1.5)
tx(s, "Enterprise security visibility and hardening — without hiring a SOC.",
   Inches(0.85), Inches(6.62), Inches(11.5), Inches(0.5), sz=15, bold=True, col=ACCENT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHO IT'S FOR  (5-second: SME SaaS healthtech logistics with lean IT)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
slide_title(s, "Who It's For", "We target deliberately — where security pressure is real and teams are lean")

# ✓ Best fit
lbox = s.shapes.add_shape(1, Inches(0.6), Inches(1.55), Inches(5.9), Inches(4.8))
lbox.fill.solid(); lbox.fill.fore_color.rgb = SURFACE; lbox.line.color.rgb = ACCENT; lbox.line.width = Pt(2)
tx(s, "✓  Best fit", Inches(0.85), Inches(1.68), Inches(5.4), Inches(0.45), sz=17, bold=True, col=ACCENT)
bul(s, [
    "Malaysian SMEs and mid-sized companies",
    "SaaS, logistics-tech, healthtech, fintech-adjacent",
    "Linux / cloud workloads (AWS, Azure, on-prem)",
    "PDPA compliance or enterprise security questionnaires",
    "Lean IT teams — no internal SOC",
], Inches(0.85), Inches(2.18), Inches(5.4), Inches(3.8), sz=14, col=TEXT, sp=8)

# ✗ Not ideal
rbox = s.shapes.add_shape(1, Inches(6.8), Inches(1.55), Inches(5.9), Inches(4.8))
rbox.fill.solid(); rbox.fill.fore_color.rgb = SURFACE; rbox.line.color.rgb = RGBColor(0x7F, 0x1D, 0x1D); rbox.line.width = Pt(2)
tx(s, "✗  Not ideal initially", Inches(7.05), Inches(1.68), Inches(5.4), Inches(0.45), sz=17, bold=True, col=RGBColor(0xFC,0x81,0x81))
bul(s, [
    "Very large enterprises needing broad coverage",
    "Banks or regulated financial institutions requiring formal SOC",
    "Expecting extensive custom consulting from day one",
], Inches(7.05), Inches(2.18), Inches(5.4), Inches(2.8), sz=14, col=MUTED, sp=8)

tx(s, "We win where security pressure is real but security ops are still lean.",
   Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.55), sz=13, italic=True, col=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HOW IT WORKS  (5-second: 4 steps, small scope, clear signal)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
slide_title(s, "How It Works", "Small scope. Clear signal. Practical output.")

steps = [
    ("01", "Onboard",  "Deploy monitoring on agreed critical systems. Define scope."),
    ("02", "Detect",   "Filter noise. Prioritize real issues. Surface what matters."),
    ("03", "Harden",   "Apply practical CIS-aligned controls on approved systems."),
    ("04", "Report",   "Clear findings, risk trends, recommended next steps monthly."),
]
sw = Inches(2.75); sg = Inches(0.38)
total = 4*sw + 3*sg; sx = (prs.slide_width - total) / 2
cols  = [ACCENT, ACCENT2, TEAL, RGBColor(0xA7,0x8A,0xF5)]
for i, (num, title, desc) in enumerate(steps):
    x = sx + i*(sw+sg)
    pill(s, num,   x, Inches(1.65), sw, Inches(0.55), bg=cols[i], sz=22)
    c = s.shapes.add_shape(1, x, Inches(2.3), sw, Inches(3.4))
    c.fill.solid(); c.fill.fore_color.rgb = SURFACE; c.line.color.rgb = cols[i]; c.line.width = Pt(1.5)
    tx(s, title, x+Inches(0.15), Inches(2.45), sw-Inches(0.3), Inches(0.5),
       sz=18, bold=True, col=cols[i], ali=PP_ALIGN.CENTER)
    tx(s, desc, x+Inches(0.15), Inches(3.05), sw-Inches(0.3), Inches(2.4),
       sz=13, col=TEXT)

tx(s, "Small scope. Clear signal. Practical improvements. Useful reporting.",
   Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.5), sz=13, italic=True, col=MUTED, ali=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PROOF  (5-second: real numbers, real findings)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
slide_title(s, "What We've Proven", "Real findings from our internal MVP demo")

stats = [
    ("2",         "Active agents\nmonitoring"),
    ("1,243",     "Events captured\nin 24 hrs"),
    ("373",       "CIS benchmark\nchecks run"),
    ("42",        "Risk score\npost-hardening"),
]
sw = Inches(2.8); sg = Inches(0.4)
total = 4*sw + 3*sg; sx = (prs.slide_width - total) / 2
cols  = [ACCENT, ACCENT2, TEAL, RGBColor(0xA7,0x8A,0xF5)]
for i, (num, desc) in enumerate(stats):
    x = sx + i*(sw+sg)
    pill(s, num, x, Inches(1.65), sw, Inches(0.95), bg=cols[i], sz=38)
    tx(s, desc, x, Inches(2.75), sw, Inches(0.75), sz=13, col=MUTED, ali=PP_ALIGN.CENTER)

divider(s, Inches(0.6), Inches(3.65), Inches(12.1), col=RGBColor(0x2A,0x35,0x4A))
tx(s, "Key findings:", Inches(0.6), Inches(3.85), Inches(12), Inches(0.4), sz=15, bold=True, col=WHITE)
bul(s, [
    "SSH brute-force detected within minutes (rule 5710)",
    "PAM auth failures flagged — rule 5503",
    "Docker DNS and service failures captured automatically",
    "373 CIS failures found; 14 remediated post-hardening",
], Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.2), sz=14, col=TEXT, sp=6)

box = s.shapes.add_shape(1, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.7))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x0D,0x28,0x1A); box.line.color.rgb = ACCENT; box.line.width = Pt(1.5)
tx(s, "Result: Ready for first customer pilot. Full service flow built and tested.",
   Inches(0.85), Inches(6.67), Inches(11.5), Inches(0.5), sz=13, bold=True, col=ACCENT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 30-DAY PILOT  (5-second: scope + outcome + price)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
slide_title(s, "Start with a 30-Day Pilot", "Narrow scope. Fast results. No long-term commitment.")

lbox = s.shapes.add_shape(1, Inches(0.6), Inches(1.55), Inches(5.9), Inches(4.0))
lbox.fill.solid(); lbox.fill.fore_color.rgb = SURFACE; lbox.line.color.rgb = ACCENT2; lbox.line.width = Pt(1.5)
tx(s, "What you get", Inches(0.85), Inches(1.68), Inches(5.4), Inches(0.4), sz=15, bold=True, col=ACCENT2)
bul(s, [
    "3–10 critical Linux servers in scope",
    "Wazuh agent deployment + enrollment",
    "Baseline findings report",
    "Practical hardening applied",
    "Management-ready final report",
], Inches(0.85), Inches(2.12), Inches(5.4), Inches(3.0), sz=14, col=TEXT, sp=8)

rbox = s.shapes.add_shape(1, Inches(6.8), Inches(1.55), Inches(5.9), Inches(4.0))
rbox.fill.solid(); rbox.fill.fore_color.rgb = SURFACE; rbox.line.color.rgb = ACCENT; rbox.line.width = Pt(1.5)
tx(s, "What you decide at Day 30", Inches(7.05), Inches(1.68), Inches(5.4), Inches(0.4), sz=15, bold=True, col=ACCENT)
bul(s, [
    "Continue on a monthly service",
    "Expand scope to more systems",
    "Or stop — no lock-in required",
], Inches(7.05), Inches(2.12), Inches(5.4), Inches(2.5), sz=14, col=TEXT, sp=8)

pbox = s.shapes.add_shape(1, Inches(0.6), Inches(5.75), Inches(12.1), Inches(1.45))
pbox.fill.solid(); pbox.fill.fore_color.rgb = SURFACE; pbox.line.color.rgb = TEAL; pbox.line.width = Pt(2)
tx(s, "Pilot pricing:", Inches(0.85), Inches(5.87), Inches(2.0), Inches(0.4), sz=13, bold=True, col=TEAL)
tx(s, "RM 2,000 – RM 5,000  (scope-dependent)", Inches(2.85), Inches(5.87), Inches(9.5), Inches(0.4), sz=18, bold=True, col=WHITE)
tx(s, "We keep the first engagement narrow on purpose — prove value fast, no transformation project required.",
   Inches(0.85), Inches(6.28), Inches(11.5), Inches(0.5), sz=12, italic=True, col=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PRICING  (5-second: 3 tiers, simple)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
slide_title(s, "Continuation After Pilot", "Start narrow. Prove value. Expand based on real need.")

tiers = [
    ("STARTER",   "RM 1,500–2,500 / mo", ACCENT,   [
        "Smaller environments",
        "Up to 5 agents",
        "Alert triage + weekly digest",
        "Monthly management report",
        "Email support",
    ]),
    ("GROWTH",    "RM 3,000–5,000 / mo", ACCENT2,  [
        "Most production teams",
        "Up to 15 agents",
        "Proactive monthly hardening",
        "Quarterly SCA scan",
        "Monthly sync call",
        "Priority support",
    ]),
    ("REGULATED", "Custom pricing",        RGBColor(0xA7,0x8A,0xF5), [
        "Stronger compliance scope",
        "Everything in Growth",
        "Weekly SCA scans",
        "Unlimited agents",
        "Dedicated analyst (PT)",
        "Incident response SLA",
    ]),
]
tw = Inches(3.6); tg = Inches(0.42)
total = 3*tw + 2*tg; sx = (prs.slide_width - total) / 2
for i, (name, price, col, items) in enumerate(tiers):
    x = sx + i*(tw+tg)
    box = s.shapes.add_shape(1, x, Inches(1.55), tw, Inches(5.0))
    box.fill.solid(); box.fill.fore_color.rgb = SURFACE
    box.line.color.rgb = col; box.line.width = Pt(2 if i==1 else 1)
    pill(s, name, x+Inches(0.4), Inches(1.7), tw-Inches(0.8), Inches(0.45), bg=col, sz=14)
    tx(s, price, x+Inches(0.15), Inches(2.25), tw-Inches(0.3), Inches(0.45),
       sz=15, bold=True, col=WHITE, ali=PP_ALIGN.CENTER)
    bul(s, items, x+Inches(0.2), Inches(2.75), tw-Inches(0.4), Inches(3.5),
       sz=12, col=TEXT, sp=6)

tx(s, "* All tiers include PDPA-aligned reporting. Custom scopes available on request.",
   Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.4),
   sz=11, italic=True, col=MUTED, ali=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CTA  (5-second: 20-min call, then pilot)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
tx(s, "Let's find out if we're a fit.",
   Inches(0.7), Inches(1.4), Inches(12), Inches(0.9), sz=40, bold=True, col=WHITE)
tx(s, "20-minute discovery call — no pitch, no pressure.",
   Inches(0.7), Inches(2.4), Inches(12), Inches(0.55), sz=20, col=MUTED)
divider(s, Inches(0.7), Inches(3.1), Inches(4), col=ACCENT)

disc = [
    ("Environment fit",     "Assess Linux/cloud workload scope"),
    ("Critical systems",     "Identify what matters most"),
    ("Visibility gaps",      "Map security blind spots"),
    ("Pilot suitability",   "Confirm fit for 30-day pilot"),
]
dw = Inches(2.75); dg = Inches(0.38)
total = 4*dw + 3*dg; sx = (prs.slide_width - total) / 2
for i, (title, desc) in enumerate(disc):
    x = sx + i*(dw+dg)
    b = s.shapes.add_shape(1, x, Inches(3.4), dw, Inches(1.8))
    b.fill.solid(); b.fill.fore_color.rgb = SURFACE; b.line.color.rgb = ACCENT2; b.line.width = Pt(1)
    tx(s, title, x+Inches(0.15), Inches(3.5), dw-Inches(0.3), Inches(0.4), sz=13, bold=True, col=ACCENT2)
    tx(s, desc, x+Inches(0.15), Inches(3.95), dw-Inches(0.3), Inches(1.0), sz=12, col=TEXT)

divider(s, Inches(0.7), Inches(5.4), Inches(12.1), col=RGBColor(0x2A,0x35,0x4A))
tx(s, "If the fit is right, Sentinel Ops can start with a limited 30-day pilot.",
   Inches(0.7), Inches(5.6), Inches(12), Inches(0.45), sz=16, bold=True, col=ACCENT)
tx(s, "[Your Name]",      Inches(0.7), Inches(6.2), Inches(12), Inches(0.45), sz=18, bold=True, col=WHITE)
tx(s, "[email]   [WhatsApp]   [LinkedIn]",
   Inches(0.7), Inches(6.6), Inches(12), Inches(0.35), sz=13, col=MUTED)

# ── Save ─────────────────────────────────────────────────────────────────────
out = "/root/.openclaw/workspace/projects/sentinel_ops/sales/Sentinel_Ops_Sales_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")